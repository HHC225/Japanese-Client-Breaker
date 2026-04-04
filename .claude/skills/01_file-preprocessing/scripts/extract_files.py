#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "openpyxl",
#     "xlrd",
#     "python-pptx",
#     "python-docx",
# ]
# ///
"""
Extract readable content from various file formats in an input directory.
Outputs a single Markdown file with all extracted content.

Run with: uv run extract_files.py <input_dir> <output_path>
uv automatically resolves dependencies declared above (PEP 723).

Supported formats:
  - Excel (.xlsx, .xls) → tables per sheet
  - PowerPoint (.pptx) → text per slide
  - Word (.docx) → full text with headings
  - CSV (.csv) → table format
  - PDF (.pdf) → metadata only (actual reading delegated to LLM Read tool)
  - Images (.png, .jpg, .jpeg, .gif, .webp) → metadata only (delegated to LLM)
  - Text (.txt, .md, .json, .xml, .yaml, .yml, .html) → direct include
"""

import sys
import os
import json
from pathlib import Path
from datetime import datetime

import openpyxl
import xlrd
from pptx import Presentation
from docx import Document


def escape_pipe(text):
    """Escape pipe characters in cell text to prevent Markdown table breakage."""
    return text.replace('|', '\\|')


def rows_to_markdown_table(rows):
    """Convert a list of row lists into a Markdown table string."""
    if not rows:
        return '(No data)\n'
    max_cols = max(len(r) for r in rows)
    rows = [r + [''] * (max_cols - len(r)) for r in rows]
    header = [escape_pipe(c) for c in rows[0]]
    table_lines = ['| ' + ' | '.join(header) + ' |']
    table_lines.append('|' + '|'.join(['---'] * max_cols) + '|')
    for row in rows[1:]:
        table_lines.append('| ' + ' | '.join(escape_pipe(c) for c in row) + ' |')
    return '\n'.join(table_lines) + '\n'


def extract_excel(filepath):
    """Extract all sheets from .xlsx Excel file as markdown tables."""
    wb = openpyxl.load_workbook(filepath, data_only=True)
    sections = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if any(cells):
                rows.append(cells)

        if not rows:
            sections.append(f"### Sheet: {sheet_name}\n\n(Empty sheet)\n")
            continue

        sections.append(f"### Sheet: {sheet_name}\n\n" + rows_to_markdown_table(rows))

    wb.close()
    return '\n'.join(sections)


def extract_excel_xls(filepath):
    """Extract all sheets from legacy .xls Excel file as markdown tables."""
    wb = xlrd.open_workbook(filepath)
    sections = []

    for sheet_name in wb.sheet_names():
        ws = wb.sheet_by_name(sheet_name)
        rows = []
        for row_idx in range(ws.nrows):
            cells = [str(ws.cell_value(row_idx, col_idx)) for col_idx in range(ws.ncols)]
            if any(c.strip() for c in cells):
                rows.append(cells)

        if not rows:
            sections.append(f"### Sheet: {sheet_name}\n\n(Empty sheet)\n")
            continue

        sections.append(f"### Sheet: {sheet_name}\n\n" + rows_to_markdown_table(rows))

    return '\n'.join(sections)


def extract_pptx(filepath):
    """Extract text from each slide in a PowerPoint file."""
    prs = Presentation(filepath)
    sections = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    texts.append('\n' + rows_to_markdown_table(rows))

        slide_content = '\n'.join(texts) if texts else '(Empty slide)'
        # Try to get slide title
        title = ''
        if slide.shapes.title and slide.shapes.title.text:
            title = f" — {slide.shapes.title.text}"

        sections.append(f"### Slide {i}{title}\n\n{slide_content}\n")

    return '\n'.join(sections)


def extract_docx(filepath):
    """Extract text from a Word document preserving heading structure."""
    doc = Document(filepath)
    sections = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style = para.style.name.lower()
        if 'heading 1' in style:
            sections.append(f"## {text}\n")
        elif 'heading 2' in style:
            sections.append(f"### {text}\n")
        elif 'heading 3' in style:
            sections.append(f"#### {text}\n")
        else:
            sections.append(f"{text}\n")

    # Extract tables
    for i, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            sections.append(f"\n**Table {i}:**\n\n" + rows_to_markdown_table(rows))

    return '\n'.join(sections)


def extract_csv(filepath):
    """Extract CSV as markdown table."""
    import csv
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return '(Empty CSV file)\n'

    return rows_to_markdown_table(rows)


def extract_tsv(filepath):
    """Extract TSV as markdown table."""
    import csv
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f, delimiter='\t')
        rows = list(reader)

    if not rows:
        return '(Empty TSV file)\n'

    return rows_to_markdown_table(rows)


def extract_text(filepath):
    """Read text file directly."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception as e:
        return f"(Error reading file: {e})\n"


def get_file_metadata(filepath):
    """Get basic file metadata."""
    stat = os.stat(filepath)
    size_kb = stat.st_size / 1024
    if size_kb > 1024:
        size_str = f"{size_kb/1024:.1f} MB"
    else:
        size_str = f"{size_kb:.1f} KB"
    modified = datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M')
    return size_str, modified


# File type handlers
EXTRACTORS = {
    '.xlsx': ('Excel Spreadsheet', extract_excel),
    '.xls': ('Excel Spreadsheet (Legacy)', extract_excel_xls),
    '.pptx': ('PowerPoint Presentation', extract_pptx),
    '.docx': ('Word Document', extract_docx),
    '.csv': ('CSV Data', extract_csv),
    '.txt': ('Text File', extract_text),
    '.md': ('Markdown File', extract_text),
    '.json': ('JSON File', extract_text),
    '.xml': ('XML File', extract_text),
    '.yaml': ('YAML File', extract_text),
    '.yml': ('YAML File', extract_text),
    '.html': ('HTML File', extract_text),
    '.htm': ('HTML File', extract_text),
    '.tsv': ('TSV Data', extract_tsv),
}

# Files that need LLM-native reading (Read tool handles these)
LLM_NATIVE = {
    '.pdf': 'PDF Document',
    '.png': 'Image (PNG)',
    '.jpg': 'Image (JPEG)',
    '.jpeg': 'Image (JPEG)',
    '.gif': 'Image (GIF)',
    '.webp': 'Image (WebP)',
    '.svg': 'Image (SVG)',
}


def process_directory(input_dir, output_path):
    """Process all files in directory and output combined markdown."""
    input_path = Path(input_dir)

    if not input_path.exists():
        print(f"Error: Input directory not found: {input_dir}", file=sys.stderr)
        sys.exit(1)

    files = sorted([f for f in input_path.iterdir() if f.is_file() and not f.name.startswith('.')])

    if not files:
        print(f"Warning: No files found in {input_dir}", file=sys.stderr)
        sys.exit(0)

    output_sections = []
    file_manifest = []
    llm_native_files = []

    output_sections.append("# Preprocessed Deliverable Content\n")
    output_sections.append(f"**Source directory**: `{input_dir}`\n")
    output_sections.append(f"**Extracted at**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
    output_sections.append(f"**Total files**: {len(files)}\n")

    # File manifest
    output_sections.append("\n## File Manifest\n")
    output_sections.append("| # | Filename | Type | Size | Status |")
    output_sections.append("|---|----------|------|------|--------|")

    for i, filepath in enumerate(files, 1):
        ext = filepath.suffix.lower()
        size_str, modified = get_file_metadata(filepath)

        if ext in EXTRACTORS:
            type_name = EXTRACTORS[ext][0]
            status = "Extracted"
        elif ext in LLM_NATIVE:
            type_name = LLM_NATIVE[ext]
            status = "LLM-Native (Read tool required)"
            llm_native_files.append((filepath, type_name))
        else:
            type_name = f"Unknown ({ext})"
            status = "Skipped"

        file_manifest.append({
            'index': i,
            'name': filepath.name,
            'type': type_name,
            'size': size_str,
            'path': str(filepath),
            'ext': ext,
            'status': status,
        })
        output_sections.append(f"| {i} | {filepath.name} | {type_name} | {size_str} | {status} |")

    output_sections.append("")

    # LLM-native files notice
    if llm_native_files:
        output_sections.append("\n## LLM-Native Files (Require Read Tool)\n")
        output_sections.append("The following files cannot be extracted by this script. The analyst agent MUST use the Read tool to access them directly:\n")
        for filepath, type_name in llm_native_files:
            output_sections.append(f"- **{filepath.name}** ({type_name}): `{filepath}`")
            if filepath.suffix.lower() == '.pdf':
                output_sections.append(f"  - Use: `Read(file_path=\"{filepath}\", pages=\"1-10\")` (adjust page range as needed)")
            else:
                output_sections.append(f"  - Use: `Read(file_path=\"{filepath}\")`")
        output_sections.append("")

    # Extract content from each supported file
    output_sections.append("\n---\n")
    output_sections.append("## Extracted Content\n")

    for filepath in files:
        ext = filepath.suffix.lower()
        if ext not in EXTRACTORS:
            continue

        type_name, extractor = EXTRACTORS[ext]
        size_str, modified = get_file_metadata(filepath)

        output_sections.append(f"\n---\n")
        output_sections.append(f"## FILE: {filepath.name}\n")
        output_sections.append(f"**Type**: {type_name} | **Size**: {size_str} | **Modified**: {modified}\n")
        output_sections.append(f"**Path**: `{filepath}`\n")

        try:
            content = extractor(str(filepath))
            if content.strip():
                output_sections.append(content)
            else:
                output_sections.append("(No content extracted)\n")
        except Exception as e:
            output_sections.append(f"**EXTRACTION ERROR**: {type(e).__name__}: {e}\n")
            print(f"Error extracting {filepath.name}: {e}", file=sys.stderr)

    # Write output
    output_text = '\n'.join(output_sections)
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    output_file.write_text(output_text, encoding='utf-8')

    # Also write manifest as JSON for programmatic access
    manifest_path = output_file.parent / '00_file_manifest.json'
    manifest_data = {
        'source_directory': str(input_dir),
        'extracted_at': datetime.now().isoformat(),
        'total_files': len(files),
        'extracted_count': sum(1 for f in file_manifest if f['status'] == 'Extracted'),
        'llm_native_count': len(llm_native_files),
        'skipped_count': sum(1 for f in file_manifest if f['status'] == 'Skipped'),
        'files': file_manifest,
        'llm_native_paths': [str(f[0]) for f in llm_native_files],
    }
    manifest_path.write_text(json.dumps(manifest_data, ensure_ascii=False, indent=2), encoding='utf-8')

    # Summary
    extracted = sum(1 for f in file_manifest if f['status'] == 'Extracted')
    native = len(llm_native_files)
    skipped = sum(1 for f in file_manifest if f['status'] == 'Skipped')

    print(f"Preprocessing complete:", file=sys.stderr)
    print(f"  Extracted: {extracted} files", file=sys.stderr)
    print(f"  LLM-Native (need Read tool): {native} files", file=sys.stderr)
    print(f"  Skipped: {skipped} files", file=sys.stderr)
    print(f"  Output: {output_path}", file=sys.stderr)
    print(f"  Manifest: {manifest_path}", file=sys.stderr)


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <input_directory> <output_markdown_path>", file=sys.stderr)
        sys.exit(1)

    process_directory(sys.argv[1], sys.argv[2])
